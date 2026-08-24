import os
import shutil
import sys
import tempfile
from pathlib import Path

import gtts
import pytest
from docx import Document
from pptx import Presentation
from pptx.util import Inches

BACKEND_DIR = Path(__file__).resolve().parents[1]
if str(BACKEND_DIR) not in sys.path:
    sys.path.insert(0, str(BACKEND_DIR))

_TEST_ROOT = Path(tempfile.mkdtemp(prefix="sahaik-tests-"))
os.environ["DATABASE_URL"] = f"sqlite:///{(_TEST_ROOT / 'suite.sqlite3').as_posix()}"
os.environ["UPLOAD_DIR"] = str(_TEST_ROOT / "uploads")


def pytest_configure(config: pytest.Config) -> None:
    config.addinivalue_line(
        "markers", "allow_llm: opt out of offline env/network stubs for this test"
    )


@pytest.fixture(autouse=True)
def _hermetic_offline(request: pytest.FixtureRequest, monkeypatch: pytest.MonkeyPatch) -> None:
    if request.node.get_closest_marker("allow_llm"):
        return
    monkeypatch.delenv("GEMINI_API_KEY", raising=False)
    monkeypatch.delenv("LLM_API_KEY", raising=False)

    def _fake_save(self: gtts.gTTS, saveaddr: object) -> None:
        Path(str(saveaddr)).write_bytes(b"ID3MOCK")

    monkeypatch.setattr(gtts.gTTS, "save", _fake_save)


def pytest_sessionfinish(session: pytest.Session, exitstatus: int) -> None:
    shutil.rmtree(_TEST_ROOT, ignore_errors=True)

KEY_PHRASE = "Photosynthesis converts light energy into chemical energy."

SAMPLE_TXT = (
    "Photosynthesis converts light energy into chemical energy. "
    "Plants absorb sunlight through chlorophyll in their leaves. "
    "The process produces glucose and releases oxygen into the air."
)

_PPTX_SENTENCES = [
    KEY_PHRASE,
    "Plants rely on chlorophyll to capture sunlight during the day.",
]


def make_txt(path: Path) -> Path:
    path.write_text(SAMPLE_TXT, encoding="utf-8")
    return path


def make_pptx(path: Path) -> Path:
    presentation = Presentation()
    for sentence in _PPTX_SENTENCES:
        slide = presentation.slides.add_slide(presentation.slide_layouts[6])
        box = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(8), Inches(2))
        box.text_frame.text = sentence
    presentation.save(str(path))
    return path


def make_docx(path: Path) -> Path:
    document = Document()
    document.add_paragraph(KEY_PHRASE)
    document.add_paragraph("Chlorophyll inside the leaves absorbs sunlight every day.")
    document.save(str(path))
    return path


def make_pdf(text: str) -> bytes:
    stream = f"BT /F1 12 Tf 72 720 Td ({text}) Tj ET".encode("latin-1")
    objects = [
        b"<< /Type /Catalog /Pages 2 0 R >>",
        b"<< /Type /Pages /Kids [3 0 R] /Count 1 >>",
        b"<< /Type /Page /Parent 2 0 R /MediaBox [0 0 612 792] "
        b"/Resources << /Font << /F1 5 0 R >> >> /Contents 4 0 R >>",
        b"<< /Length " + str(len(stream)).encode("latin-1") + b" >>\nstream\n"
        + stream + b"\nendstream",
        b"<< /Type /Font /Subtype /Type1 /BaseFont /Helvetica >>",
    ]
    out = bytearray(b"%PDF-1.4\n")
    offsets: list[int] = []
    for number, body in enumerate(objects, start=1):
        offsets.append(len(out))
        out += f"{number} 0 obj\n".encode("latin-1") + body + b"\nendobj\n"
    xref_offset = len(out)
    out += f"xref\n0 {len(objects) + 1}\n".encode("latin-1")
    out += b"0000000000 65535 f \n"
    for offset in offsets:
        out += f"{offset:010d} 00000 n \n".encode("latin-1")
    out += (
        f"trailer\n<< /Size {len(objects) + 1} /Root 1 0 R >>\n"
        f"startxref\n{xref_offset}\n%%EOF\n"
    ).encode("latin-1")
    return bytes(out)


@pytest.fixture
def txt_file(tmp_path: Path) -> Path:
    return make_txt(tmp_path / "sample.txt")


@pytest.fixture
def pptx_file(tmp_path: Path) -> Path:
    return make_pptx(tmp_path / "sample.pptx")


@pytest.fixture
def docx_file(tmp_path: Path) -> Path:
    return make_docx(tmp_path / "sample.docx")


@pytest.fixture
def pdf_data() -> bytes:
    return make_pdf(KEY_PHRASE)
