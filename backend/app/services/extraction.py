from io import BytesIO

import docx
import pptx
import re
from pypdf import PdfReader

_BOILERPLATE_LINE = re.compile(
    r"(?i)(?:copyright|\u00a9|all rights reserved|permission is (?:hereby )?granted|"
    r"provided proper attribution|reproduce (?:portions of )?this (?:material|publication|paper)|"
    r"specially granted|journalistic|scholarly works|request[s]? to (?:pubs|permissions)|"
    r"liability|warranty|creative commons|arxiv:\d{4}\.\d+|under review|preprint)"
)


def _normalize(text: str) -> str:
    normalized = text.replace("\r\n", "\n").replace("\r", "\n")
    kept = [
        line
        for line in normalized.splitlines()
        if not (_BOILERPLATE_LINE.search(line) and len(line.split()) < 60)
    ]
    normalized = "\n".join(kept)
    while "\n\n\n" in normalized:
        normalized = normalized.replace("\n\n\n", "\n\n")
    return normalized.strip()


def _extension(filename: str) -> str:
    return filename.lower().rsplit(".", 1)[-1] if "." in filename else ""


def extract_text(filename: str, data: bytes) -> str:
    """Dispatch by extension: .pptx .pdf .docx .txt. Raise ValueError on unsupported/empty."""
    ext = _extension(filename)
    if ext == "pptx":
        presentation = pptx.Presentation(BytesIO(data))
        raw = "\n".join(
            shape.text_frame.text
            for slide in presentation.slides
            for shape in slide.shapes
            if shape.has_text_frame
        )
    elif ext == "pdf":
        reader = PdfReader(BytesIO(data))
        raw = "\n".join(page.extract_text() or "" for page in reader.pages)
    elif ext == "docx":
        document = docx.Document(BytesIO(data))
        raw = "\n".join(paragraph.text for paragraph in document.paragraphs)
    elif ext == "txt":
        raw = data.decode("utf-8", errors="replace")
    else:
        raise ValueError(f"unsupported file type: {filename}")
    result = _normalize(raw)
    if len(result) < 10:
        raise ValueError("empty document")
    return result


def doc_type(filename: str) -> str:
    return _extension(filename)
